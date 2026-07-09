import re
from pathlib import Path

PROJECT = Path(__file__).parent.name


def get_images_dir(agent):
    """
    Resolve images directory via agent's data/config context.
    Images must end up under www/{project}/images/ so gateway serves them.
    Strategy: find 'www' directory relative to plugins, falling back to agent paths.
    """
    # Try relative to this plugin file: plugins/applications/{project}/ → ../../../www/{project}/images
    candidate = Path(__file__).parent.parent.parent / "www" / PROJECT / "images"
    if candidate.parent.exists():
        candidate.mkdir(parents=True, exist_ok=True)
        return candidate

    # Fallback: use agent.data_dir and hope www is sibling
    data_root = Path(agent.data_dir).parent
    candidate2 = data_root / "www" / PROJECT / "images"
    candidate2.mkdir(parents=True, exist_ok=True)
    return candidate2


# ==================== LLM Cleanup (shared: web + PDF) ====================

MAX_LLM_CHARS = 6000

NOISE_LINE_PATTERNS = [
    r'^\s*责任编辑[：:].*$',
    r'^\s*海量资讯[、，].*$',
    r'^\s*(VIP课程|APP专享).*$',
    r'^\s*加载中\.\.\.\s*$',
    r'^\s*(上一页|下一页)\s*$',
    r'^\s*\d+/\d+\s*$',
    r'^\s*(分享|举报|收藏|点赞|评论\s*\d*|阅读\s*\d+)\s*$',
    r'^\s*(相关阅读|热门推荐|猜你喜欢|为你推荐|点击加载更多|展开全文|阅读全文)\s*$',
    r'^\s*(7X24小时|交易提示|操盘必读|股市直播|财经头条|商品行情|外汇计算器|基金净值|最新公告|限售解禁|数据中心|条件选股|千股千评|个股诊断|大宗交易|业绩预告).*$',
    r'^\s*(徐小明|凯恩斯|占豪|花荣|wu2198|叶檀|曹中铭|股民大张|宇辉战舰|杨伟民|温彬|余华莘|李德林|李庚南|程实).*$',
    r'^\s*\d{2}/\d{2}\s.*$',
    r'^\s*\d{2}-\d{2}\s.*$',
    r'^\s*来源[：:]\s*\S+$',
    r'^\s*原特斯拉.*$',
    r'^\s*重磅.*$',
    r'^\s*BJ30.*$',
    r'^\s*韩泰.*$',
    r'^\s*央视曝光.*$',
    r'^\s*2026\s*款.*$',
    r'^\s*王朝网.*$',
    r'^\s*DeepMind.*$',
    r'^\s*乘联分会.*$',
    r'^\s*月薪从.*$',
    r'^\s*-\s*\d+\s*-\s*$',
    r'^\s*\d+\s*$',
]


def compress_text(raw_text):
    lines = raw_text.split('\n')
    kept = []

    for line in lines:
        stripped = line.strip()
        if not stripped:
            kept.append('')
            continue

        if any(re.match(p, stripped, re.I) for p in NOISE_LINE_PATTERNS):
            continue

        if len(stripped) < 15 and not stripped.startswith('#') and not stripped.startswith('-'):
            if not re.match(r'^[A-Z\u4e00-\u9fff][A-Za-z\u4e00-\u9fff\s]{3,}$', stripped):
                continue

        kept.append(stripped)

    result = '\n'.join(kept)
    result = re.sub(r'\n{3,}', '\n\n', result)
    result = result.strip()

    if len(result) > MAX_LLM_CHARS:
        split_point = int(MAX_LLM_CHARS * 0.75)
        head = result[:split_point]
        tail = result[split_point:MAX_LLM_CHARS]
        tail_lines = tail.split('\n')
        tail_kept = [l for l in tail_lines if len(l.strip()) > 40 or l.strip() == '']
        result = head + '\n' + '\n'.join(tail_kept)

    return result[:MAX_LLM_CHARS]


LLM_CLEANUP_PROMPT = """Clean up this raw text into well-structured Markdown.

Rules:
- Remove ALL noise: ads, navigation, "related articles", "hot topics", stock info, author lists, pagination, page numbers, repeated headers, watermark text, app download prompts.
- Pay attention to the END — website footers and PDF appendices often contain unrelated links or reference lists. DELETE them.
- KEEP: title, headings, body paragraphs, quotes, image markers [Image: ...], table content.
- Structure: # for title, ## for sections, ### for subsections, - for bullets, > for quotes, | tables | for tabular data.
- Do NOT summarize. Preserve all factual content. Just remove garbage.
- Output ONLY Markdown. No explanations.

Raw text:
---
{raw_text}
---"""


async def ai_clean_text(agent, raw_text, source_name="document"):
    llm = agent.llm
    if not llm:
        return raw_text

    compressed = compress_text(raw_text)
    agent.log.info(f"AI clean ({source_name}): {len(raw_text)} → {len(compressed)} chars")

    try:
        prompt = LLM_CLEANUP_PROMPT.format(raw_text=compressed)
        result = await llm.chat([{"role": "user", "content": prompt}])
        agent.log.info(f"AI clean done: {len(result)} chars")
        return result.strip()
    except Exception as e:
        agent.log.warning(f"AI clean failed for {source_name}: {e}")
        return raw_text


# ==================== DOCX ====================

def extract_image_from_docx(docx, rels, r_id, images_dir, counter):
    if r_id not in rels:
        return ""
    rel = rels[r_id]
    if "image" not in rel.reltype:
        return ""
    try:
        image_part = rel.target_part
        ext = Path(image_part.partname).suffix or ".png"
        fname = f"docx_img_{counter[0]}{ext}"
        counter[0] += 1
        fpath = images_dir / fname
        fpath.write_bytes(image_part.blob)
        return f"![{fname}](images/{fname})"
    except Exception:
        return ""


def parse_docx_paragraph(para, docx, images_dir, img_counter):
    from docx.oxml.ns import qn

    parts = []
    for run in para.runs:
        text = run.text
        if not text:
            continue
        t = text
        if run.bold:
            t = f"**{t}**"
        if run.italic:
            t = f"*{t}*"
        if run.underline:
            t = f"<u>{t}</u>"

        drawings = run._element.findall(qn('w:drawing'))
        for drawing in drawings:
            blip = drawing.findall('.//' + qn('a:blip'))
            for b in blip:
                embed = b.get(qn('r:embed'))
                if embed:
                    md_img = extract_image_from_docx(docx, para.part.rels, embed, images_dir, img_counter)
                    parts.append("\n" + md_img + "\n")

        parts.append(t)

    return "".join(parts)


def convert_docx(file_path, images_dir):
    from docx import Document
    from docx.oxml.ns import qn

    doc = Document(file_path)
    img_counter = [0]
    md_lines = []

    for para in doc.paragraphs:
        style = para.style.name if para.style else ""

        drawings = para._element.findall(qn('w:drawing'))
        img_md_parts = []
        for d in drawings:
            blips = d.findall('.//' + qn('a:blip'))
            for blip in blips:
                embed = blip.get(qn('r:embed'))
                if embed:
                    md_img = extract_image_from_docx(doc, para.part.rels, embed, images_dir, img_counter)
                    if md_img:
                        img_md_parts.append(md_img)

        if img_md_parts and not para.text.strip():
            md_lines.extend(img_md_parts)
            continue

        text = parse_docx_paragraph(para, doc, images_dir, img_counter)
        if not text.strip():
            if img_md_parts:
                md_lines.extend(img_md_parts)
            else:
                md_lines.append("")
            continue

        line = ""
        if style.startswith("Heading 1") or style.startswith("heading 1"):
            line = f"# {text}"
        elif style.startswith("Heading 2") or style.startswith("heading 2"):
            line = f"## {text}"
        elif style.startswith("Heading 3") or style.startswith("heading 3"):
            line = f"### {text}"
        elif style.startswith("Heading 4") or style.startswith("heading 4"):
            line = f"#### {text}"
        elif "List Bullet" in style or "List Paragraph" in style:
            num_id = para._element.find(qn('w:numPr'))
            level = 0
            if num_id is not None:
                ilvl = num_id.find(qn('w:ilvl'))
                if ilvl is not None:
                    level = int(ilvl.get(qn('w:val'), 0))
            indent = "  " * level
            line = f"{indent}- {text}"
        elif "List Number" in style:
            line = f"1. {text}"
        else:
            line = text

        if img_md_parts:
            line += "\n" + "\n".join(img_md_parts)
        md_lines.append(line)

    for table in doc.tables:
        md_lines.append("")
        for i, row in enumerate(table.rows):
            cells = []
            for cell in row.cells:
                cell_text = cell.text.replace("\n", " ").replace("|", "\\|")
                cells.append(cell_text)
            md_lines.append("| " + " | ".join(cells) + " |")
            if i == 0:
                md_lines.append("| " + " | ".join(["---"] * len(cells)) + " |")
        md_lines.append("")

    return "\n".join(md_lines)


# ==================== XLSX ====================

def is_header_row(row_data, prev_row_data=None):
    if not row_data:
        return False
    non_empty = [c for c in row_data if c]
    if not non_empty:
        return False
    str_count = sum(1 for c in non_empty if not str(c).replace('.', '').replace('-', '').isdigit())
    num_count = len(non_empty) - str_count
    avg_len = sum(len(str(c)) for c in non_empty) / len(non_empty)
    if str_count >= num_count and avg_len < 40:
        return True
    if prev_row_data:
        prev_types = [type(c).__name__ for c in prev_row_data if c]
        curr_types = [type(c).__name__ for c in row_data if c]
        if prev_types != curr_types:
            return True
    return False


def find_table_regions(rows):
    regions = []
    current_region = []
    for row in rows:
        non_empty = [c for c in row if c is not None and str(c).strip()]
        if not non_empty:
            if current_region:
                regions.append(current_region)
                current_region = []
        else:
            current_region.append(row)
    if current_region:
        regions.append(current_region)
    return regions


def convert_xlsx(file_path, images_dir):
    from openpyxl import load_workbook

    wb = load_workbook(file_path, data_only=True)
    img_counter = [0]
    md_lines = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        md_lines.append(f"## {sheet_name}")
        md_lines.append("")

        if hasattr(ws, '_images') and ws._images:
            for img in ws._images:
                ext = Path(img.path).suffix if hasattr(img, 'path') and img.path else ".png"
                fname = f"xlsx_img_{img_counter[0]}{ext}"
                img_counter[0] += 1
                fpath = images_dir / fname
                fpath.write_bytes(img._data())
                md_lines.append(f"![{fname}](images/{fname})")
            md_lines.append("")

        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            md_lines.append("*Empty sheet*")
            md_lines.append("")
            continue

        regions = find_table_regions(rows)
        if not regions:
            md_lines.append("*Empty sheet*")
            md_lines.append("")
            continue

        for region_idx, region in enumerate(regions):
            if len(regions) > 1:
                md_lines.append(f"### {sheet_name} — Region {region_idx + 1}")
                md_lines.append("")

            max_cols = max(len(row) for row in region)
            has_header = is_header_row(region[0]) if region else False

            for i, row in enumerate(region):
                cells = []
                for j in range(max_cols):
                    val = row[j] if j < len(row) else ""
                    if val is None:
                        val = ""
                    val = str(val).replace("\n", " ").replace("|", "\\|")
                    cells.append(val)
                md_lines.append("| " + " | ".join(cells) + " |")
                if i == 0 and has_header:
                    md_lines.append("| " + " | ".join(["---"] * max_cols) + " |")

            md_lines.append("")

    return "\n".join(md_lines)


# ==================== PPTX ====================

def is_title_shape(shape):
    try:
        if shape.is_placeholder:
            ph = shape.placeholder_format
            if ph.type == 1:
                return 'title'
            if ph.type == 2:
                return 'body'
            if ph.type == 3:
                return 'subtitle'
            if ph.type == 7:
                return 'title'
    except Exception:
        pass
    return None


def get_shape_text_hierarchy(shape):
    items = []
    if not shape.has_text_frame:
        return items
    for para in shape.text_frame.paragraphs:
        text = para.text.strip()
        if text:
            items.append((para.level if para.level else 0, text))
    return items


def convert_pptx(file_path, images_dir):
    from pptx import Presentation

    prs = Presentation(file_path)
    img_counter = [0]
    md_lines = []

    for slide_num, slide in enumerate(prs.slides, 1):
        has_content = False

        title_text = ""
        if slide.shapes.title:
            title_text = slide.shapes.title.text.strip()
            if title_text:
                md_lines.append(f"## {title_text}")
                md_lines.append("")
                has_content = True

        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            placeholder_type = is_title_shape(shape)
            if placeholder_type == 'subtitle':
                items = get_shape_text_hierarchy(shape)
                for level, text in items:
                    md_lines.append(f"### {text}")
                    md_lines.append("")
                    has_content = True

        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            placeholder_type = is_title_shape(shape)
            if placeholder_type == 'subtitle':
                continue

            if shape.has_text_frame and placeholder_type != 'title':
                items = get_shape_text_hierarchy(shape)
                if items:
                    first_level, first_text = items[0]
                    if placeholder_type == 'body' and first_level == 0 and len(items) > 1:
                        md_lines.append(f"### {first_text}")
                        md_lines.append("")
                        items = items[1:]

                    for level, text in items:
                        indent = "  " * level
                        md_lines.append(f"{indent}- {text}")
                    md_lines.append("")
                    has_content = True

            if shape.has_table:
                table = shape.table
                md_lines.append("")
                for i, row in enumerate(table.rows):
                    cells = [cell.text.replace("\n", " ").replace("|", "\\|") for cell in row.cells]
                    md_lines.append("| " + " | ".join(cells) + " |")
                    if i == 0:
                        md_lines.append("| " + " | ".join(["---"] * len(cells)) + " |")
                md_lines.append("")
                has_content = True

            if shape.shape_type == 13:
                try:
                    image = shape.image
                    ext = image.content_type.split("/")[-1] if "/" in image.content_type else "png"
                    fname = f"pptx_img_{img_counter[0]}.{ext}"
                    img_counter[0] += 1
                    fpath = images_dir / fname
                    fpath.write_bytes(image.blob)
                    md_lines.append(f"![{fname}](images/{fname})")
                    md_lines.append("")
                    has_content = True
                except Exception:
                    pass

            if shape.has_chart:
                try:
                    chart = shape.chart
                    md_lines.append("*Chart:*")
                    if chart.has_title:
                        md_lines.append(f"**{chart.chart_title.text_frame.text}**")
                        md_lines.append("")
                    for series in chart.series:
                        md_lines.append(f"- {series.name if hasattr(series, 'name') else 'Series'}")
                    md_lines.append("")
                    has_content = True
                except Exception:
                    md_lines.append("*[Chart]*")
                    md_lines.append("")

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                md_lines.append("> **Notes:**")
                for line in notes.split("\n"):
                    md_lines.append(f"> {line}")
                md_lines.append("")
                has_content = True

        if has_content:
            md_lines.append("---")
            md_lines.append("")

    return "\n".join(md_lines)


# ==================== PDF ====================

def convert_pdf(file_path, images_dir):
    import fitz

    doc = fitz.open(file_path)
    img_counter = [0]
    lines = []

    for page_num in range(len(doc)):
        page = doc[page_num]

        text = page.get_text("text")
        if text.strip():
            lines.append(text.strip())

        for img_index, img in enumerate(page.get_images(full=True)):
            try:
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                ext = base_image["ext"]
                fname = f"pdf_img_{img_counter[0]}.{ext}"
                img_counter[0] += 1
                fpath = images_dir / fname
                fpath.write_bytes(image_bytes)
                lines.append(f"[Image: {fname}]")
            except Exception:
                pass

    doc.close()

    raw_text = '\n\n'.join(lines)

    try:
        import pdfplumber
        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages):
                tables = page.extract_tables()
                if tables:
                    raw_text += "\n\n"
                    for table in tables:
                        for row in table:
                            if row and any(c for c in row if c):
                                cells = [str(c).replace("\n", " ").replace("|", "\\|") if c else "" for c in row]
                                raw_text += "| " + " | ".join(cells) + " |\n"
                        raw_text += "\n"
    except Exception:
        pass

    return raw_text


# ==================== Main Entry ====================

async def execute(envelop, agent):
    payload = envelop.payload.get("payload", envelop.payload)
    file_path = payload.get("file_path", "")

    if not file_path:
        envelop.payload = {"error": "file_path required"}
        return envelop

    p = Path(file_path)
    if not p.exists():
        envelop.payload = {"error": f"File not found: {file_path}"}
        return envelop

    suffix = p.suffix.lower()

    images_dir = get_images_dir(agent)
    images_dir.mkdir(parents=True, exist_ok=True)

    for f in images_dir.glob("*"):
        try:
            f.unlink()
        except Exception:
            pass

    try:
        if suffix == ".docx":
            md_content = convert_docx(file_path, images_dir)
            file_type = "Word Document"
        elif suffix == ".xlsx":
            md_content = convert_xlsx(file_path, images_dir)
            file_type = "Excel Spreadsheet"
        elif suffix == ".pptx":
            md_content = convert_pptx(file_path, images_dir)
            file_type = "PowerPoint Presentation"
        elif suffix == ".pdf":
            raw_text = convert_pdf(file_path, images_dir)
            file_type = "PDF Document"
            if agent.llm:
                md_content = await ai_clean_text(agent, raw_text, "PDF")
            else:
                md_content = raw_text
        else:
            envelop.payload = {"error": f"Unsupported format: {suffix}"}
            return envelop
    except Exception as e:
        agent.log.error(f"Conversion error: {e}")
        envelop.payload = {"error": f"Conversion failed: {str(e)}"}
        return envelop

    md_dir = Path(agent.data_dir) / PROJECT
    md_dir.mkdir(parents=True, exist_ok=True)
    md_path = md_dir / f"{p.stem}_converted.md"
    md_path.write_text(md_content, encoding="utf-8")

    images = [f.name for f in images_dir.glob("*")] if images_dir.exists() else []

    envelop.payload = {
        "ok": True,
        "md_content": md_content,
        "file_type": file_type,
        "original_name": p.name,
        "md_path": str(md_path),
        "images": images,
        "images_dir": str(images_dir),
    }
    return envelop

